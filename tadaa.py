import os
import utils
import ppt
from uuid import uuid4
import datetime
from pptx import Presentation
from main import UPLOAD_DIR, ROOT_PATH
import pandas as pd


def generate_ppt(form_data:dict, export_files:list, schema:dict):
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    target_files = get_target_files_from_schema(schema)
    PROJECT_DIR = create_project_dir(timestamp)
    save_csvs(export_files, PROJECT_DIR)
    parsed_export_files = parse_export_files(export_files, target_files, PROJECT_DIR)
    tadaabject = create_tadaabject(form_data, parsed_export_files, timestamp)
    ppt_path = setup_powerpoint(schema, tadaabject['audit_data'], PROJECT_DIR)
    tadaabject['ppt_path'] = ppt_path

    return tadaabject

def create_project_dir(timestamp):
    PROJECT_DIR = UPLOAD_DIR + '/' + timestamp
    os.makedirs(PROJECT_DIR, exist_ok=True)
    return PROJECT_DIR

def save_csvs(export_files, dir):
    for file in export_files:
        file.save(os.path.join(dir, os.path.basename(file.filename)))

# TODO: Do I need to parse more in order to populate the powerpoint?
def parse_export_files(export_files:list, target_files:list, PROJECT_DIR:str):
    matched_export_files = [file for file in export_files if any(target_file in file.filename for target_file in target_files)]
    
    parsed_export_files = []
    for file in matched_export_files:
        file_name = os.path.basename(file.filename)
        filepath = PROJECT_DIR + '/' + file_name
        dataframe = pd.read_csv(filepath)

        file_object = {
            "name": file_name,
            "num": dataframe.shape[0],
            #"data": dataframe
        }
        parsed_export_files.append(file_object)

    return parsed_export_files


def create_tadaabject(form_data:dict, parsed_export_files:dict, timestamp:str):
    audit_data = (form_data, parsed_export_files)

    tadaabject = {
        "domain": form_data["domain_url"],
        "ts": timestamp,
        "audit_data": audit_data,
        "ppt_path": ""
    }
    
    return tadaabject


def setup_powerpoint(schema:dict, audit_data:tuple, PROJECT_DIR:str):
    form_data = audit_data[0]
    
    template_name = schema['ppt']
    template_path = ROOT_PATH + '/ppts/pptx/' + template_name
    presentation = Presentation(template_path)

    populated_presentation = populate_powerpoint(presentation, schema, audit_data)

    save_path = PROJECT_DIR + f"/{form_data['domain_url']}.pptx"
    populated_presentation.save(save_path)

    return save_path

# TODO: Figure out how to populate powerpoint
def populate_powerpoint(presentation:Presentation, schema:dict, audit_data:tuple):
    schema_shapes = get_shape_ids_from_schema(schema)
    template_slides = [slide for slide in presentation.slides]

    for slide in template_slides:
        for template_shape in slide.shapes:
            if template_shape.has_text_frame:
                for shape in schema_shapes:
                    shape_id, shape_index = shape
                    if shape_id == template_shape.name:
                        # TODO: insert data into ppt
                        # insert_data_into_ppt(data, type, shape)
                        pass
    
    return presentation

def get_shape_ids_from_schema(schema:dict):  # Also gets the slide index for shape (id, index)
    shapes = []
    for slide in schema["slides"]:
        if "keys" in slide and "shapes" in slide["keys"]:
            shape_id = slide["keys"]["shapes"]["id"][0]
            shape_index = slide['index']
            shapes.append((shape_id, shape_index))

    return shapes
    

def get_target_files_from_schema(schema:dict):
    target_files = []
    for slide in schema["slides"]:
        if "target_export_files" in slide:
            target_files.extend(slide["target_export_files"])
    target_files = [file for file in target_files if file]

    return target_files

# TODO: Variable names from exports match the target file name. Variable names from form match input id.
# TODO: Add input id field to schema?

# Slide 7.
# ga_bool = form_data['ga_bool'] 
# sc_bool = form_data['sc_bool'] 

# # Slide 9.
# sc_mob_usability = form_data['sc_mob_usability'] 
# num_mob_friendly_issues = form_data['num_mob_friendly_issues']

# # Slide 10.
# is_sitemap_submitted_sc = form_data['is_sitemap_submitted_sc']  
# sitemap_slug = form_data['sitemap-slug']
# url_is_orphaned_and_was_not_found_by_the_crawler = self.count() 
# redirect__3xx__url_in_xml_sitemaps = self.count()
# noindex_url_in_xml_sitemaps = self.count()
# urls_in_multiple_sitemaps = self.count()

# # Slide 11.
# has_robots = form_data['has_robots']
# robots_blocks_good_pages = form_data['robots_blocks_good_pages']
# robots_should_block_bad_pages = form_data['robots_should_block_bad_pages']

# # Slide 12.
# has_structured_data_errors = form_data['has_structured_data_errors']
# num_structured_data_errors = form_data['num_structured_data_errors']

# # Slide 14.
# title_tag_length_too_long = self.count()
# title_tag_length_too_short = self.count()
# urls_with_duplicate_page_titles = self.count()

# # Slide 15.
# description_length_too_long = self.count()
# description_length_too_short = self.count()
# urls_with_duplicate_meta_descriptions = self.count()
# meta_description_is_missing = self.count()
# meta_description_is_empty = self.count()

# # Slide 16.
# urls_with_duplicate_h1s = self.count()
# h1__tag_is_empty = self.count()

# # Slide 17.
# has_duplicate_content = form_data['has_duplicate_content']
# has_thin_content = form_data['has_thin_content']

# # Slide 19.
# has_good_cta = form_data['has_good_cta']

# # Slide 20.
# has_onsite_blog = form_data['has_onsite_blog']
# is_blog_updated = form_data['is_blog_updated']

# # Slide 21.
# images_with_missing_alt_text = self.count()

# # Slide 23.
# broken_internal_urls = self.count()
# broken_external_urls = self.count()
# external_url_redirect_broken__4xx_or_5xx = self.count()
# sc_404s = form_data['sc_404s']
# sc_crawl_anomalies = form_data['sc_crawl_anomalies']
# sc_soft_404s = form_data['sc_soft_404s']

# # Slide 24.
# has_canonicals = form_data['has_canonicals']
# canonical_has_errors = form_data['canonical_has_errors']

# # Slide 25.
# internal_redirected_urls = self.count()
# external_redirected_urls = self.count()

# # Slide 26.
# loads_http = form_data['loads_http']
# loads_mixed_resources = form_data['loads_mixed_resources']
# ssl_exp = form_data['ssl_exp']

# # Slide 27.
# mobile_load_time = form_data['mobile_load_time']
# desktop_load_time = form_data['desktop_load_time']

# # Slide 30.
# broken_backlinks = form_data['broken_backlinks']




# ======== OLD BELOW =========




def old_generate_ppt(project_dir:str, form_data:dict, root_path:str, project_name:str, timestamp:str, schema):
    """
    Generates a populated Powerpoint document using the custom data object.
    
    @param project_dir: Project path by timestamp "os.path/uploads/[timestamp]"
    @param manual_data: Data submitted through form fields
    @param root_path: the path to the root of application

    @return [dict] tadaabject: Our custom tech audit data object
    """
    # Form the tadaabject
    audits_id = str(uuid4())
    client_id = str(uuid4())
    domain = form_data["domain_url"]
    project_type = "InvalidType"

    tadaabject  = {
        "audits_id": audits_id,
        "client_id": client_id,
        "client_name": project_name,
        "domain": domain,
        "ts": timestamp,
        "project_type": project_type,
        #"ppt_url": ppt_path, # This gives a valid server path that is useless client-side
        "audit_data": ""
    }
    old_parse_data(project_dir, form_data)
    
    tadaabject['ppt_url'] = ppt._populate_powerpoint(schema, data)
    
    return tadaabject


def old_parse_data(project_dir, form_data):
    """
    Takes the upload directory and organizes the data necessary for the tech audit. Returns our custom data object.

    @param [str] upload_dir: path to the server's upload directory.

    @return [dict] final_data_obj: our custom data object {target_file: [data]} that contains Pandas dataframe/series objects.
    """

    all_uploaded_files = os.listdir(project_dir)

    matched_list = utils.match_target_hint_files(all_uploaded_files)
    matched_paths = utils.get_abs_paths(matched_list, project_dir)

    # All the CSV data
    final_data_obj = utils.get_data_obj(matched_paths)

    return None
    #return final_data_obj