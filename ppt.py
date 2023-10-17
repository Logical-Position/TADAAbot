from pptx import Presentation

# Slide 7.
has_sc_access = bool
has_ga_access = bool

# Slide 9.
is_mobile_friendly = bool
num_mob_friendly_issues = 0

# Slide 10.
is_sitemap_submitted_sc = bool
sitemap_errors = 0
sitemap_location = "sitemap_location"
orphaned_pages = 0
redirects_in_sitemap = 0
noindex_urls_in_sitemap = 0
urls_in_multiple_sitemaps = 0

# Slide 11.
has_robots = bool
blocks_good_pages = bool
should_block_bad_pages = bool
robots_location = "robots_location"

# Slide 12.
has_str_data = bool
num_str_data_err = 0

# Slide 14.
titles_too_long = 0
titles_too_short = 0
duplicate_titles = 0

# Slide 15.
desc_too_long = 0
desc_too_short = 0
duplicate_desc = 0
desc_missing = 0
desc_empty = 0

# Slide 16.
duplicate_h1_tags = 0
h1_tag_empty = 0

# Slide 17.
has_duplicate_content = bool
has_thin_content = bool

# Slide 19.
is_good_cta = bool

# Slide 20.
has_onsite_blog = bool
is_blog_updated = bool

# Slide 21.
img_alt_text = 0

# Slide 23.
broken_int_links = 0
broken_ext_links = 0
broken_4xx_5xx = 0
sc_404s = 0
sc_crawl_anom = 0
sc_soft_404s = 0

# Slide 24.
has_canonicals = bool
canonical_has_errors = bool

# Slide 25.
int_redirect_links = 0
ext_redirect_links = 0

# Slide 26.
loads_http = bool
loads_mixed_resources = bool
ssl_exp = bool

# Slide 27.
mob_load_time = float
desk_load_time = float

# Slide 30.
broken_backlinks = 0


def populate_powerpoint(final_data_object, project_dir, root_path, project_name):
    """
    Goes through the template PowerPoint slide by slide, and adjusts the values/SEO recommendation text to correspond
    to the calculated data and user inputs.
    @param exports_path: the path to the main exports folder.
    @param root_path: the path to the tech audit template ppt.
    """
    template_path = root_path + '/SEOC Tech Audit Template.pptx'
    presentation = Presentation(template_path)
    slides = [slide for slide in presentation.slides]
    slide_num = 0
    for slide in slides:
        slide_num += 1
        print(f'---------Slide {slide_num}---------')
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraph = text_frame.paragraphs[0]
                runs = paragraph.runs

                # Slide 7
                if 'ga_bool' == shape.name:
                    if not has_ga_access:
                        runs[0].text = 'No'
                    print(shape.name)
                if 'sc_bool' == shape.name:
                    if not has_sc_access:
                        runs[0].text = 'No'
                    print(shape.name)

                # Slide 9
                if 'sc_mob_usability_analyst_notes' == shape.name:
                    if not is_mobile_friendly:
                        runs[0].text = f"We found that {num_mob_friendly_issues} of your pages are NOT mobile-friendly. We recommend making your pages mobile-friendly, so visitors can easily view your site on their mobile devices."
                    print(shape.name)
                if 'sc_mob_usability' == shape.name:
                    print(shape.name)

                # Slide 10
                if 'sitemap_analyst_notes' == shape.name:
                    print(shape.name)
                    runs[0].text = f'Sitemap Errors: {sitemap_errors} | Orphaned Pages: {orphaned_pages} | ' \
                                   f'Redirects in Sitemap: {redirects_in_sitemap} | Noindex URLs in Sitemap: ' \
                                   f'{noindex_urls_in_sitemap} | URLs in Multiple Sitemaps: {urls_in_multiple_sitemaps}'
                if 'sc_sitemap_submitted_bool' == shape.name:
                    if not is_sitemap_submitted_sc:
                        runs[0].text = "We could not find a sitemap. We recommend creating a sitemap and submitting it through Google Search Console."
                    print(shape.name)
                if 'sitemap_errors_bool' == shape.name:
                    print(shape.name)
                    if sitemap_errors == 0:
                        runs[0].text = 'No'

                # Slide 11
                if 'robots_analyst_notes' == shape.name:
                    if not has_robots:
                        runs[0].text = "We found that your site does not have a robots.txt file. We recommend creating a robots.txt file, so you can block pages that do not need to be indexed."   
                    else:
                        if blocks_good_pages:
                            runs[0].text = 'We found that your site has pages that should be indexed being blocked by noindex tags. We recommend removing the noindex tags to prevent this.'
                        if should_block_bad_pages:
                            runs[0].text = 'We found that your site has a robots.txt file; however, it does not block shopping cart pages, pages that contain login information, and other account pages. We recommend updating the robots.txt file to block these pages.'
                        else:
                            runs[0].text += (f"Your robots.txt is located here: " + robots_location)

                    print(shape.name)

                # Slide 12
                if 'structured_data_analyst_notes' == shape.name:
                    if not has_str_data:
                        runs[0].text = "We found that your website does not have structured data. Structured data can be beneficial to you and your customers. Accurate structured data can help with rankings, while visitors immediately receive pertinent information about the website. We recommend creating structured data for your website."
                    elif num_str_data_err > 0:
                        runs[0].text = f"We found that your website has structured data; however, there are {num_str_data_err} items with errors. We recommend updating the missing fields or descriptions to eliminate these errors."
                    print(shape.name)

                # Slide 14
                if 'meta_title_analyst_notes' == shape.name:
                    print(shape.name)
                    runs[0].text = f'Short titles : {titles_too_short} | Long titles :' \
                                   f' {titles_too_long} | Duplicate titles: {duplicate_titles}'

                # Slide 15
                if 'meta_descriptions_analyst_notes' == shape.name:
                    print(shape.name)
                    runs[0].text = f'Missing descriptions: {desc_missing} | Short descriptions: {desc_too_short} | ' \
                                   f'Long descriptions: {desc_too_long} | Empty descriptions: {desc_empty}'

                # Slide 16
                if 'h1_analyst_notes' == shape.name:
                    print(shape.name)
                    runs[0].text = f'Empty H1 tags: {h1_tag_empty} | Duplicate H1 tags: {duplicate_h1_tags}'

                # Slide 17
                if 'site_content_analyst_notes' == shape.name:
                    if not has_thin_content:
                        runs[0].text = f"After reviewing your site, we found that you have quite a bit of content. While this is a great start, we recommend adding keyword-targeted content. This will help crawlers understand the products and services your website offers, and you will also be able to rank better in the long run."
                    print(shape.name)

                # Slide 18
                if 'dup_content_analyst_notes' == shape.name:
                    if has_duplicate_content:
                        runs[0].text = "We found that there are several instances of duplicate content across your site, primarily on your category pages. We recommend writing original content for each page. This can increase readability and rankings."
                    print(shape.name)

                # Slide 19
                if 'cta_analyst_notes' == shape.name:
                    if not is_good_cta:
                        runs[0].text = "We found that your site could benefit from visible CTAs. Clear contact forms, phone numbers, and addresses can help with click-through-rate and conversions. We recommend creating accurate and visible CTAs throughout your website."
                    print(shape.name)

                # Slide 20
                if 'blog_updated_reg_bool' == shape.name:
                    if not is_blog_updated:
                        runs[0].text = "No"
                    print(shape.name)
                if 'onsite_blog_bool' == shape.name:
                    if not has_onsite_blog:
                        runs[0].text = "No"
                    print(shape.name)

                # Slide 21
                if 'alt_text_analyst_notes' == shape.name:
                    print(shape.name)
                    runs[0].text = f"We found that {img_alt_text} of your website's images do not have image alt-text. We recommend adding alt-text to your images. This will give a keyword-focused image description to users and search engines. You can view the images that do not have alt-text here."
                    if img_alt_text == 0:
                        runs[0].text = "We found that your website's images have image alt-text."

                # Slide 23
                if '404s_analyst_notes' == shape.name:
                    print(shape.name)

                    if broken_4xx_5xx or broken_ext_links or broken_int_links > 0:

                        runs[0].text = f"We found that your site has {broken_int_links} broken internal links and {broken_ext_links + broken_4xx_5xx} broken external links. You can find the broken internal links here, and the broken external links here. Google Search Console also found # 404 errors; they can be found here. However, we identified that only # need to be addressed. # crawl anomalies were also found, they can be found here. It also found # soft 404 errors, which can be found here. A soft 404 is a URL that returns a page telling the user that the page does not exist and also a 200-level (success) code. In some cases, it might be a page with little or no content--for example, a sparsely populated or empty page. If there is a new URL for any of these 404 errors, we recommend redirecting the broken URL to that new page. If no URL exists, we recommend redirecting the broken URL to a relevant 2xx page. We recommend replacing the broken links with a new link or removing them from any pages they appear on."

                # Slide 24
                if 'canonical_analyst_notes' == shape.name:
                    if not has_canonicals and canonical_has_errors:
                        runs[0].text = "We found that your canonicals are not set up properly. We recommend pointing filtered or paginated pages back to the first page to avoid duplicate content. This includes tagged collection pages, which are dynamically generated and cannot be edited. We recommend pointing tagged pages back to the original collection page."
                    elif not has_canonicals:
                        runs[0].text = "We found that your pages do not have canonicals; however, it looks like canonicals are unnecessary at the moment since you do not have a lot of duplicate content. If you add a blog, canonicals may be necessary."
                        if canonical_has_errors:
                            runs[0].text = "We found that your pages do not have canonicals. We recommend pointing filtered or paginated pages back to the first page to avoid duplicate content."
                    print(shape.name)

                # Slide 25
                if 'redirects_analyst_notes' == shape.name:
                    if int_redirect_links > 0 or ext_redirect_links > 0:
                        runs[0].text = f"We found {int_redirect_links} internal redirect links and {ext_redirect_links} external redirect links on your site. You can find the internal redirect links here, and the external redirect links here. We recommend correcting these links with their new 2XX-status URLs."
                    
                    print(shape.name)


                # Slide 26
                if 'site_security_analyst_notes' == shape.name:
                    if loads_http:
                        runs[0].text = "We found that your site is not secure, as it does not load in HTTPS. We recommend that your site switch to HTTPS.\n"
                    if loads_mixed_resources and not loads_http:
                        runs[0].text += "We found that your site has HTTPS; however, it loads mixed resources. We recommend ensuring that your site's resources load securely.\n"
                    if ssl_exp and not loads_http:
                        runs[0].text += "We found that your site has HTTPS; however, the SSL certificate has expired. We recommend renewing your SSL certificate.\n"
                    
                    print(shape.name)

                # Slide 27
                if 'desk_load_time_float' == shape.name:
                    runs[0].text = str(desk_load_time)
                    print(shape.name)
                if 'mob_load_time_float' == shape.name:
                    runs[0].text = str(mob_load_time)
                    print(shape.name)

                # Slide 29
                if 'dom_auth_analyst_notes' == shape.name:
                    print(shape.name)

                # Slide 30
                if 'backlinks_analyst_notes' == shape.name:
                    if broken_backlinks > 0:
                        runs[0].text = f"We found {broken_backlinks} broken backlinks. You can view them here. To help maintain your website's authority, we recommend fixing broken backlinks from relevant, high-authority domains."
                    print(shape.name)

    output_path = project_dir + f'/{project_name}.pptx'
    presentation.save(output_path)

    return output_path
