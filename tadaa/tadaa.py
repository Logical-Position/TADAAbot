'''
Action:
POST /gen-ppt
1. Get CSV filenames
2. Generate PowerPoint
    1. Create timestamp
    2. Select target CSVs
    3. Create a project directory
    4. Save CSVs
    5. Parse CSVs
    6. Create tadaabject
    7. Setup PowerPoint
    8. Add PowerPoint path to tadaabject
    9. Return tadaabject
3. Return tadaabject to client

Rework to:
1. Create UUID
1. Save CSVs to Uploads
2. Save images to Uploads
3. Generate PowerPoint
    1. Run Audit
        2. Select target CSVs (can this be done client-side?)
        3. Parse CSVs
        4. Parse other form data
        5. Return audit_data
    2. Create PowerPoint
        1. Setup PowerPoint
        2. Populate PowerPoint
    3. Return tadaabject
4. Return tadaabject
'''

from tadaa import csv
from tadaa import ppt

def run_audit(sitebulb_files, target_export_files, form_data):
    sitebulb_data = csv.parse_sitebulb_csvs(sitebulb_files, target_export_files)
    print('')
    print('=== FORM DATA ===')
    print(form_data)
    print('')
    audit_data = {}
    audit_data.update(sitebulb_data)
    audit_data.update(form_data)
    # parse_sc_404_files
    # parse_siteliner_results
    # parse ...
    # concatenate all the data
    print('')
    print('=== AUDIT DATA ===')
    print(audit_data)
    print('')
    return audit_data

def create_presentation(template_path, output_path, audit_data, schema):
    ppt.create_powerpoint(template_path, output_path, audit_data, schema)









































'''

import os
from tadaa import csv
from tadaa import ppt

from pptx import Presentation
import pandas as pd



def get_audit_data():
    return {}

def generate_ppt(form_data:dict, export_files:list, schema:dict):
    
    ppt.old_populate_powerpoint()

    return {}

def new_generate_ppt(form_data:dict, export_files:list, schema:dict, upload_dir):
    target_files = get_target_files_from_schema(schema)

    parsed_export_files = parse_export_files(export_files, target_files, PROJECT_DIR)
    tadaabject = create_tadaabject(form_data, parsed_export_files, timestamp)
    
    ppt_path = setup_powerpoint(schema, tadaabject['audit_data'], PROJECT_DIR)
    tadaabject['ppt_path'] = ppt_path

    return tadaabject

# TODO: Do I need to parse more in order to populate the powerpoint?
def parse_export_files(export_files:list, target_files:list, PROJECT_DIR:str):
    matched_export_files = [file for file in export_files if any(target_file in file.filename for target_file in target_files)]
    
    parsed_export_files = []
    for file in matched_export_files:
        file_name = os.path.basename(file.filename)
        filepath = PROJECT_DIR + '/' + file_name
        dataframe = pd.read_csv(filepath)

        file_object = {
            "name": file_name,
            "num": dataframe.shape[0],
            #"data": dataframe
        }
        parsed_export_files.append(file_object)

    return parsed_export_files


def create_tadaabject(form_data:dict, parsed_export_files:dict, timestamp:str):
    audit_data = (form_data, parsed_export_files)

    tadaabject = {
        "domain": form_data["domain_url"],
        "ts": timestamp,
        "audit_data": audit_data,
        "ppt_path": ""
    }
    
    return tadaabject


def setup_powerpoint(schema:dict, audit_data:tuple, PROJECT_DIR:str):
    form_data = audit_data[0]
    
    template_name = schema['ppt']
    template_path = '/ppts/pptx/' + template_name
    print(" TEMPLATE PATH")
    print(template_path)
    presentation = Presentation(template_path)

    populated_presentation = populate_powerpoint(presentation, schema, audit_data)

    save_path = PROJECT_DIR + f"/{form_data['domain_url']}.pptx"
    populated_presentation.save(save_path)

    return save_path

# TODO: Figure out how to populate powerpoint
def populate_powerpoint(presentation:Presentation, schema:dict, audit_data:tuple):
    schema_shapes = get_shape_ids_from_schema(schema)
    template_slides = [slide for slide in presentation.slides]

    for slide in template_slides:
        for template_shape in slide.shapes:
            if template_shape.has_text_frame:
                for shape in schema_shapes:
                    shape_id, shape_index = shape
                    if shape_id == template_shape.name:
                        # TODO: insert data into ppt
                        # insert_data_into_ppt(data, type, shape)
                        pass
    
    return presentation

def get_shape_ids_from_schema(schema:dict):  # Also gets the slide index for shape (id, index)
    shapes = []
    for slide in schema["slides"]:
        if "keys" in slide and "shapes" in slide["keys"]:
            shape_id = slide["keys"]["shapes"]["id"][0]
            shape_index = slide['index']
            shapes.append((shape_id, shape_index))

    return shapes
    

def get_target_files_from_schema(schema:dict):
    target_files = []
    for slide in schema["slides"]:
        if "target_export_files" in slide:
            target_files.extend(slide["target_export_files"])
    target_files = [file for file in target_files if file]

    return target_files

# ======== OLD BELOW =========

def old_generate_ppt(project_dir:str, form_data:dict, root_path:str, project_name:str, timestamp:str, schema):
    """
    Generates a populated Powerpoint document using the custom data object.
    
    @param project_dir: Project path by timestamp "os.path/uploads/[timestamp]"
    @param manual_data: Data submitted through form fields
    @param root_path: the path to the root of application

    @return [dict] tadaabject: Our custom tech audit data object
    """

    audits_id = str(uuid4())
    client_id = str(uuid4())
    domain = manual_data["domain_url"]
    project_type = "InvalidType"

    # Not doing anything with parsed_data yet
    parsed_data = parse_data(project_dir, manual_data)
    pop_ppt = ppt.populate_powerpoint(parsed_data, project_dir, root_path, project_name)

    tadaabject  = {
        "audits_id": audits_id,
        "client_id": client_id,
        "client_name": project_name,
        "domain": domain,
        "ts": timestamp,
        "project_type": project_type,
        "ppt_url": pop_ppt,
    }
    old_parse_data(project_dir, form_data)
    
    tadaabject['ppt_url'] = ppt._populate_powerpoint(schema, {})
    
    return tadaabject



def parse_data(project_dir, manual_data):
    """
    Takes the upload directory and organizes the data necessary for the tech audit. Returns our custom data object.

    @param [str] upload_dir: path to the server's upload directory.

    @return [dict] final_data_obj: our custom data object {target_file: [data]} that contains Pandas dataframe/series objects.
    """

    all_uploaded_files = os.listdir(project_dir)

    matched_list = csv.match_target_hint_files(all_uploaded_files)
    matched_paths = csv.get_abs_paths(matched_list, project_dir)

    # All the CSV data
    final_data_obj = csv.get_data_obj(matched_paths)

    return None
    #return final_data_obj
'''
